# SPDX-License-Identifier: MIT
"""
M1 Image Acquisition Service for Multimodal Knowledge Graph System

Handles image acquisition from:
- Documents (PDF, DOCX, PPTX)
- Web URLs
- File uploads
- External APIs
"""

import os
import io
import hashlib
import mimetypes
import re
import json
from datetime import datetime
from pathlib import Path
from typing import List, Optional, Dict, Any, Union, Tuple
from dataclasses import dataclass, field, asdict
from enum import Enum
import urllib.request
import urllib.robotparser
import ssl
import base64

from PIL import Image
from loguru import logger
from pydantic import BaseModel

# Try importing document processing libraries
try:
    import fitz  # PyMuPDF for PDF
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False
    logger.warning("PyMuPDF not available - PDF extraction disabled")

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available - DOCX extraction disabled")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available - PPTX extraction disabled")

try:
    import requests
    REQUESTS_AVAILABLE = True
except ImportError:
    REQUESTS_AVAILABLE = False
    logger.warning("requests not available - web crawling will use urllib")

from app.core.config import settings


class ImageType(str, Enum):
    """Image type classification"""
    FACILITY = "facility"
    FLOOR_PLAN = "floor_plan"
    SCENE = "scene"
    STANDARD = "standard"
    SIGN = "sign"
    UNKNOWN = "unknown"


class SourceType(str, Enum):
    """Image source type"""
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    WEB = "web"
    UPLOAD = "upload"
    API = "api"


@dataclass
class ImageData:
    """Data class for image information"""
    file_path: str
    thumbnail_path: Optional[str] = None
    source: str = ""
    source_type: str = SourceType.UPLOAD.value
    source_context: Optional[str] = None
    original_width: int = 0
    original_height: int = 0
    file_size: int = 0
    image_type: str = ImageType.UNKNOWN.value
    tags: List[str] = field(default_factory=list)
    image_id: Optional[str] = None
    metadata: Dict[str, Any] = field(default_factory=dict)
    created_at: str = field(default_factory=lambda: datetime.utcnow().isoformat())
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary"""
        return asdict(self)
    
    def generate_id(self) -> str:
        """Generate unique ID for the image"""
        if self.image_id is None:
            hash_input = f"{self.file_path}_{self.created_at}_{self.file_size}"
            self.image_id = hashlib.md5(hash_input.encode()).hexdigest()[:16]
        return self.image_id


class DocumentImageExtractor:
    """Extract images from PDF, DOCX, and PPTX documents"""
    
    def __init__(self, raw_path: Optional[Path] = None, thumbnail_path: Optional[Path] = None):
        self.raw_path = raw_path or settings.images_raw
        self.thumbnail_path = thumbnail_path or settings.images_thumbnails
        self.raw_path.mkdir(parents=True, exist_ok=True)
        self.thumbnail_path.mkdir(parents=True, exist_ok=True)
        
        # Keyword patterns for image classification
        self.classification_patterns = {
            ImageType.FLOOR_PLAN: [
                r'floor', r'plan', r'layout', r'blueprint', r'arrangement',
                r'平面图', r'布局', r'设计图'
            ],
            ImageType.FACILITY: [
                r'facility', r'equipment', r'machine', r'system', r'plant',
                r'设备', r'设施', r'机器', r'装置'
            ],
            ImageType.SCENE: [
                r'scene', r'site', r'location', r'photo', r'view', r'landscape',
                r'场景', r'现场', r'现场照片', r'风景'
            ],
            ImageType.SIGN: [
                r'sign', r'label', r'symbol', r'warning', r'safety', r'mark',
                r'标志', r'标识', r'符号', r'警示', r'安全'
            ],
            ImageType.STANDARD: [
                r'standard', r'spec', r'specification', r'diagram', r'schematic',
                r'标准', r'规范', r'规格', r'原理图'
            ]
        }
    
    def extract_from_pdf(self, pdf_path: Union[str, Path]) -> List[ImageData]:
        """Extract images from PDF file"""
        if not PYPDF2_AVAILABLE:
            logger.error("PyMuPDF not available for PDF extraction")
            return []
        
        images = []
        pdf_path = Path(pdf_path)
        
        if not pdf_path.exists():
            logger.error(f"PDF file not found: {pdf_path}")
            return []
        
        try:
            doc = fitz.open(str(pdf_path))
            doc_name = pdf_path.stem
            
            for page_num, page in enumerate(doc):
                image_list = page.get_images(full=True)
                
                for img_idx, img_info in enumerate(image_list):
                    try:
                        xref = img_info[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]
                        
                        # Generate filename
                        filename = f"{doc_name}_p{page_num+1}_img{img_idx+1}.{image_ext}"
                        file_path = self.raw_path / filename
                        
                        # Save image
                        with open(file_path, "wb") as f:
                            f.write(image_bytes)
                        
                        # Create ImageData
                        image_data = self._create_image_data(
                            file_path=file_path,
                            source=str(pdf_path),
                            source_type=SourceType.PDF.value,
                            source_context=f"PDF page {page_num + 1}, image {img_idx + 1}"
                        )
                        images.append(image_data)
                        
                    except Exception as e:
                        logger.warning(f"Failed to extract image {img_idx} from page {page_num}: {e}")
            
            doc.close()
            logger.info(f"Extracted {len(images)} images from PDF: {pdf_path.name}")
            
        except Exception as e:
            logger.error(f"Error extracting images from PDF {pdf_path}: {e}")
        
        return images
    
    def extract_from_docx(self, docx_path: Union[str, Path]) -> List[ImageData]:
        """Extract images from DOCX file"""
        if not DOCX_AVAILABLE:
            logger.error("python-docx not available for DOCX extraction")
            return []
        
        images = []
        docx_path = Path(docx_path)
        
        if not docx_path.exists():
            logger.error(f"DOCX file not found: {docx_path}")
            return []
        
        try:
            doc = DocxDocument(str(docx_path))
            doc_name = docx_path.stem
            
            for rel_idx, rel in enumerate(doc.part.rels.values()):
                if "image" in rel.target_ref:
                    try:
                        image_part = rel.target_part.blob
                        image_ext = rel.target_ref.split('.')[-1]
                        
                        # Generate filename
                        filename = f"{doc_name}_img{rel_idx+1}.{image_ext}"
                        file_path = self.raw_path / filename
                        
                        # Save image
                        with open(file_path, "wb") as f:
                            f.write(image_part)
                        
                        # Create ImageData
                        image_data = self._create_image_data(
                            file_path=file_path,
                            source=str(docx_path),
                            source_type=SourceType.DOCX.value,
                            source_context=f"DOCX image {rel_idx + 1}"
                        )
                        images.append(image_data)
                        
                    except Exception as e:
                        logger.warning(f"Failed to extract image {rel_idx}: {e}")
            
            logger.info(f"Extracted {len(images)} images from DOCX: {docx_path.name}")
            
        except Exception as e:
            logger.error(f"Error extracting images from DOCX {docx_path}: {e}")
        
        return images
    
    def extract_from_pptx(self, pptx_path: Union[str, Path]) -> List[ImageData]:
        """Extract images from PPTX file"""
        if not PPTX_AVAILABLE:
            logger.error("python-pptx not available for PPTX extraction")
            return []
        
        images = []
        pptx_path = Path(pptx_path)
        
        if not pptx_path.exists():
            logger.error(f"PPTX file not found: {pptx_path}")
            return []
        
        try:
            prs = Presentation(str(pptx_path))
            ppt_name = pptx_path.stem
            
            for slide_idx, slide in enumerate(prs.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, "image"):
                        try:
                            image_blob = shape.image.blob
                            image_ext = shape.image.ext
                            
                            # Generate filename
                            filename = f"{ppt_name}_s{slide_idx+1}_img{shape_idx+1}.{image_ext}"
                            file_path = self.raw_path / filename
                            
                            # Save image
                            with open(file_path, "wb") as f:
                                f.write(image_blob)
                            
                            # Create ImageData
                            image_data = self._create_image_data(
                                file_path=file_path,
                                source=str(pptx_path),
                                source_type=SourceType.PPTX.value,
                                source_context=f"PPTX slide {slide_idx + 1}, shape {shape_idx + 1}"
                            )
                            images.append(image_data)
                            
                        except Exception as e:
                            logger.warning(f"Failed to extract shape image: {e}")
            
            logger.info(f"Extracted {len(images)} images from PPTX: {pptx_path.name}")
            
        except Exception as e:
            logger.error(f"Error extracting images from PPTX {pptx_path}: {e}")
        
        return images
    
    def _create_image_data(self, file_path: Path, source: str, 
                          source_type: str, source_context: str) -> ImageData:
        """Create ImageData from file path with classification"""
        try:
            with Image.open(file_path) as img:
                width, height = img.size
                img_format = img.format.lower() if img.format else "unknown"
            
            file_size = file_path.stat().st_size
            
            # Classify image type
            image_type = self._classify_image_type(str(file_path), source_context)
            
            # Generate tags based on classification
            tags = self._generate_tags(image_type, source_context)
            
            return ImageData(
                file_path=str(file_path),
                source=source,
                source_type=source_type,
                source_context=source_context,
                original_width=width,
                original_height=height,
                file_size=file_size,
                image_type=image_type,
                tags=tags
            )
            
        except Exception as e:
            logger.warning(f"Error creating ImageData for {file_path}: {e}")
            return ImageData(
                file_path=str(file_path),
                source=source,
                source_type=source_type,
                source_context=source_context
            )
    
    def _classify_image_type(self, file_path: str, context: str = "") -> str:
        """Classify image type based on filename and context"""
        combined_text = f"{file_path} {context}".lower()
        
        # Score each category
        scores = {}
        for img_type, patterns in self.classification_patterns.items():
            score = 0
            for pattern in patterns:
                if re.search(pattern, combined_text, re.IGNORECASE):
                    score += 1
            scores[img_type] = score
        
        # Return highest scoring type, or unknown
        if max(scores.values()) > 0:
            return max(scores, key=scores.get).value
        return ImageType.UNKNOWN.value
    
    def _generate_tags(self, image_type: str, context: str = "") -> List[str]:
        """Generate tags for image"""
        tags = [image_type]
        if context:
            # Extract key terms from context
            words = re.findall(r'\b[a-zA-Z]+\b', context.lower())
            tags.extend(words[:5])  # Add first 5 words as tags
        return list(set(tags))


class WebImageCrawler:
    """Crawl and extract images from web URLs"""
    
    def __init__(self, raw_path: Optional[Path] = None, thumbnail_path: Optional[Path] = None):
        self.raw_path = raw_path or settings.images_raw
        self.thumbnail_path = thumbnail_path or settings.images_thumbnails
        self.raw_path.mkdir(parents=True, exist_ok=True)
        self.thumbnail_path.mkdir(parents=True, exist_ok=True)
        
        self.robot_parser = urllib.robotparser.RobotFileParser()
        self.session_headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        }
        
        # Classification patterns (reused)
        self.classification_patterns = {
            ImageType.FLOOR_PLAN: [r'floor', r'plan', r'layout', r'blueprint'],
            ImageType.FACILITY: [r'facility', r'equipment', r'machine', r'system'],
            ImageType.SCENE: [r'scene', r'site', r'photo', r'view'],
            ImageType.SIGN: [r'sign', r'label', r'symbol', r'warning'],
            ImageType.STANDARD: [r'standard', r'spec', r'diagram', r'schematic']
        }
    
    def crawl_images(self, url: str, keywords: List[str] = None, 
                    max_images: int = 50) -> List[ImageData]:
        """Crawl images from URL with keyword filtering"""
        images = []
        
        try:
            # Check robots.txt
            if not self._can_fetch(url):
                logger.warning(f"Disallowed by robots.txt: {url}")
                return []
            
            # Fetch page content
            if REQUESTS_AVAILABLE:
                response = requests.get(url, headers=self.session_headers, timeout=10)
                content = response.content
            else:
                req = urllib.request.Request(url, headers=self.session_headers)
                with urllib.request.urlopen(req, timeout=10) as resp:
                    content = resp.read()
            
            # Extract image URLs
            image_urls = self._extract_image_urls(content, url)
            
            # Filter by keywords
            if keywords:
                image_urls = self._filter_by_keywords(image_urls, keywords)
            
            # Limit to max_images
            image_urls = image_urls[:max_images]
            
            # Download and save images
            for img_url in image_urls:
                img_data = self._download_image(img_url, url)
                if img_data:
                    images.append(img_data)
            
            logger.info(f"Crawled {len(images)} images from {url}")
            
        except Exception as e:
            logger.error(f"Error crawling {url}: {e}")
        
        return images
    
    def _can_fetch(self, url: str) -> bool:
        """Check if URL is allowed by robots.txt"""
        try:
            parsed_url = urllib.parse.urlparse(url)
            robots_url = f"{parsed_url.scheme}://{parsed_url.netloc}/robots.txt"
            
            self.robot_parser.set_url(robots_url)
            self.robot_parser.read()
            
            return self.robot_parser.can_fetch("*", url)
        except Exception as e:
            logger.warning(f"Could not check robots.txt for {url}: {e}")
            return True  # Allow by default if can't check
    
    def _extract_image_urls(self, content: bytes, base_url: str) -> List[str]:
        """Extract image URLs from HTML content"""
        import re
        import urllib.parse
        
        urls = []
        
        try:
            html = content.decode('utf-8', errors='ignore')
            
            # Find img tags
            img_patterns = [
                r'<img[^>]+src=["\']([^"\']+)["\']',
                r'<image[^>]+href=["\']([^"\']+)["\']',
                r'url\(["\']?([^"\'()]+)["\']?\)'
            ]
            
            for pattern in img_patterns:
                matches = re.findall(pattern, html, re.IGNORECASE)
                for match in matches:
                    if not match.startswith(('data:', 'javascript:', 'mailto:')):
                        # Resolve relative URLs
                        full_url = urllib.parse.urljoin(base_url, match)
                        if self._is_valid_image_url(full_url):
                            urls.append(full_url)
            
        except Exception as e:
            logger.warning(f"Error extracting image URLs: {e}")
        
        return list(set(urls))
    
    def _is_valid_image_url(self, url: str) -> bool:
        """Check if URL points to a valid image"""
        image_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.svg')
        return any(url.lower().endswith(ext) for ext in image_extensions)
    
    def _filter_by_keywords(self, urls: List[str], keywords: List[str]) -> List[str]:
        """Filter URLs by keywords"""
        filtered = []
        keywords_lower = [k.lower() for k in keywords]
        
        for url in urls:
            url_lower = url.lower()
            # Check if any keyword is in the URL
            if any(kw in url_lower for kw in keywords_lower):
                filtered.append(url)
            else:
                # Include if no keywords specified (conservative approach)
                filtered.append(url)
        
        return filtered
    
    def _download_image(self, img_url: str, source_url: str) -> Optional[ImageData]:
        """Download single image and create ImageData"""
        try:
            if REQUESTS_AVAILABLE:
                response = requests.get(img_url, headers=self.session_headers, timeout=10)
                content = response.content
            else:
                req = urllib.request.Request(img_url, headers=self.session_headers)
                with urllib.request.urlopen(req, timeout=10) as resp:
                    content = resp.read()
            
            # Generate filename
            parsed = urllib.parse.urlparse(img_url)
            filename = os.path.basename(parsed.path) or f"web_{hashlib.md5(img_url.encode()).hexdigest()[:8]}.jpg"
            file_path = self.raw_path / filename
            
            # Save image
            with open(file_path, "wb") as f:
                f.write(content)
            
            # Create ImageData
            image_data = self._create_image_data(file_path, source_url, img_url)
            return image_data
            
        except Exception as e:
            logger.warning(f"Failed to download image {img_url}: {e}")
            return None
    
    def _create_image_data(self, file_path: Path, source_url: str, 
                          image_url: str) -> ImageData:
        """Create ImageData with classification"""
        try:
            with Image.open(file_path) as img:
                width, height = img.size
            
            file_size = file_path.stat().st_size
            
            # Classify based on URL
            image_type = self._classify_from_url(image_url)
            
            return ImageData(
                file_path=str(file_path),
                source=source_url,
                source_type=SourceType.WEB.value,
                source_context=image_url,
                original_width=width,
                original_height=height,
                file_size=file_size,
                image_type=image_type,
                tags=[image_type]
            )
            
        except Exception as e:
            logger.warning(f"Error creating ImageData for {file_path}: {e}")
            return ImageData(
                file_path=str(file_path),
                source=source_url,
                source_type=SourceType.WEB.value,
                source_context=image_url
            )
    
    def _classify_from_url(self, url: str) -> str:
        """Classify image type based on URL"""
        url_lower = url.lower()
        
        scores = {}
        for img_type, patterns in self.classification_patterns.items():
            score = sum(1 for p in patterns if p in url_lower)
            scores[img_type] = score
        
        if max(scores.values()) > 0:
            return max(scores, key=scores.get).value
        return ImageType.UNKNOWN.value


class FileUploadHandler:
    """Handle uploaded image files"""
    
    def __init__(self, raw_path: Optional[Path] = None, thumbnail_path: Optional[Path] = None):
        self.raw_path = raw_path or settings.images_raw
        self.thumbnail_path = thumbnail_path or settings.images_thumbnails
        self.raw_path.mkdir(parents=True, exist_ok=True)
        self.thumbnail_path.mkdir(parents=True, exist_ok=True)
        
        self.thumbnail_size = (256, 256)
    
    def handle_upload(self, file_data: bytes, filename: str, 
                     metadata: Optional[Dict[str, Any]] = None) -> ImageData:
        """Process uploaded file and create ImageData"""
        try:
            # Generate unique filename
            timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
            safe_filename = re.sub(r'[^\w\-.]', '_', filename)
            new_filename = f"upload_{timestamp}_{safe_filename}"
            file_path = self.raw_path / new_filename
            
            # Save original file
            with open(file_path, "wb") as f:
                f.write(file_data)
            
            # Extract metadata
            img_metadata = self._extract_metadata(file_path)
            
            # Generate thumbnail
            thumbnail_path = self._generate_thumbnail(file_path)
            
            # Classify image
            image_type = self._classify_image(file_path, metadata)
            
            # Create tags
            tags = self._generate_tags(image_type, metadata)
            
            return ImageData(
                file_path=str(file_path),
                thumbnail_path=str(thumbnail_path) if thumbnail_path else None,
                source="upload",
                source_type=SourceType.UPLOAD.value,
                original_width=img_metadata.get("width", 0),
                original_height=img_metadata.get("height", 0),
                file_size=len(file_data),
                image_type=image_type,
                tags=tags,
                metadata=metadata or {}
            )
            
        except Exception as e:
            logger.error(f"Error handling upload {filename}: {e}")
            raise
    
    def _extract_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract basic image metadata"""
        try:
            with Image.open(file_path) as img:
                width, height = img.size
                img_format = img.format
                mode = img.mode
            
            return {
                "width": width,
                "height": height,
                "format": img_format,
                "mode": mode
            }
        except Exception as e:
            logger.warning(f"Error extracting metadata: {e}")
            return {}
    
    def _generate_thumbnail(self, file_path: Path) -> Optional[Path]:
        """Generate thumbnail for image"""
        try:
            with Image.open(file_path) as img:
                # Convert to RGB if necessary
                if img.mode in ('RGBA', 'P'):
                    img = img.convert('RGB')
                
                # Create thumbnail
                img.thumbnail(self.thumbnail_size, Image.Resampling.LANCZOS)
                
                # Save thumbnail
                thumb_filename = f"thumb_{file_path.name}"
                thumb_path = self.thumbnail_path / thumb_filename
                img.save(thumb_path, "JPEG", quality=85)
                
                return thumb_path
                
        except Exception as e:
            logger.warning(f"Error generating thumbnail for {file_path}: {e}")
            return None
    
    def _classify_image(self, file_path: Path, metadata: Optional[Dict[str, Any]] = None) -> str:
        """Classify image type"""
        # Simple heuristic based on filename and size
        filename_lower = file_path.stem.lower()
        
        if any(kw in filename_lower for kw in ['floor', 'plan', 'layout']):
            return ImageType.FLOOR_PLAN.value
        elif any(kw in filename_lower for kw in ['facility', 'equipment', 'machine']):
            return ImageType.FACILITY.value
        elif any(kw in filename_lower for kw in ['sign', 'label', 'mark']):
            return ImageType.SIGN.value
        elif any(kw in filename_lower for kw in ['standard', 'spec', 'diagram']):
            return ImageType.STANDARD.value
        elif metadata and 'type' in metadata:
            return metadata['type']
        
        return ImageType.UNKNOWN.value
    
    def _generate_tags(self, image_type: str, metadata: Optional[Dict[str, Any]] = None) -> List[str]:
        """Generate tags for image"""
        tags = [image_type, "upload"]
        
        if metadata:
            if 'tags' in metadata:
                tags.extend(metadata['tags'])
            if 'description' in metadata:
                words = re.findall(r'\b\w+\b', metadata['description'].lower())
                tags.extend(words[:10])
        
        return list(set(tags))


class ImageAcquisitionService:
    """Main service for image acquisition from multiple sources"""
    
    def __init__(self):
        self.raw_path = settings.images_raw
        self.thumbnail_path = settings.images_thumbnails
        self.raw_path.mkdir(parents=True, exist_ok=True)
        self.thumbnail_path.mkdir(parents=True, exist_ok=True)
        
        # Initialize sub-services
        self.document_extractor = DocumentImageExtractor(
            raw_path=self.raw_path,
            thumbnail_path=self.thumbnail_path
        )
        self.web_crawler = WebImageCrawler(
            raw_path=self.raw_path,
            thumbnail_path=self.thumbnail_path
        )
        self.upload_handler = FileUploadHandler(
            raw_path=self.raw_path,
            thumbnail_path=self.thumbnail_path
        )
        
        logger.info("ImageAcquisitionService initialized")
    
    def acquire_from_document(self, doc_path: Union[str, Path]) -> List[str]:
        """Acquire images from document (PDF, DOCX, PPTX)"""
        doc_path = Path(doc_path)
        image_ids = []
        
        if not doc_path.exists():
            logger.error(f"Document not found: {doc_path}")
            return []
        
        suffix = doc_path.suffix.lower()
        
        try:
            if suffix == '.pdf':
                images = self.document_extractor.extract_from_pdf(doc_path)
            elif suffix == '.docx':
                images = self.document_extractor.extract_from_docx(doc_path)
            elif suffix == '.pptx':
                images = self.document_extractor.extract_from_pptx(doc_path)
            else:
                logger.error(f"Unsupported document format: {suffix}")
                return []
            
            # Generate thumbnails and IDs
            for img in images:
                self._generate_thumbnail(img.file_path)
                image_ids.append(img.generate_id())
            
            logger.info(f"Acquired {len(image_ids)} images from document: {doc_path.name}")
            
        except Exception as e:
            logger.error(f"Error acquiring from document {doc_path}: {e}")
        
        return image_ids
    
    def acquire_from_url(self, url: str, keywords: List[str] = None) -> List[str]:
        """Acquire images from web URL"""
        image_ids = []
        
        try:
            images = self.web_crawler.crawl_images(url, keywords)
            
            for img in images:
                self._generate_thumbnail(img.file_path)
                image_ids.append(img.generate_id())
            
            logger.info(f"Acquired {len(image_ids)} images from URL: {url}")
            
        except Exception as e:
            logger.error(f"Error acquiring from URL {url}: {e}")
        
        return image_ids
    
    def acquire_from_upload(self, file_data: bytes, filename: str,
                          metadata: Optional[Dict[str, Any]] = None) -> str:
        """Acquire single image from file upload"""
        try:
            img_data = self.upload_handler.handle_upload(file_data, filename, metadata)
            return img_data.generate_id()
            
        except Exception as e:
            logger.error(f"Error acquiring from upload: {e}")
            raise
    
    def acquire_from_api(self, api_endpoint: str, params: Dict[str, Any] = None) -> List[str]:
        """Acquire images from external API"""
        image_ids = []
        
        try:
            # Make API request
            if REQUESTS_AVAILABLE:
                response = requests.get(api_endpoint, params=params, timeout=30)
                response.raise_for_status()
                data = response.json()
            else:
                import json as json_module
                req = urllib.request.Request(api_endpoint + '?' + urllib.parse.urlencode(params or {}))
                with urllib.request.urlopen(req, timeout=30) as resp:
                    data = json_module.load(resp)
            
            # Handle different response formats
            image_urls = self._parse_api_response(data)
            
            # Download each image
            for img_url in image_urls:
                try:
                    if REQUESTS_AVAILABLE:
                        img_response = requests.get(img_url, timeout=10)
                        content = img_response.content
                    else:
                        req = urllib.request.Request(img_url)
                        with urllib.request.urlopen(req, timeout=10) as resp:
                            content = resp.read()
                    
                    # Save and process
                    filename = f"api_{hashlib.md5(img_url.encode()).hexdigest()[:8]}.jpg"
                    file_path = self.raw_path / filename
                    with open(file_path, "wb") as f:
                        f.write(content)
                    
                    self._generate_thumbnail(str(file_path))
                    
                    img_data = ImageData(
                        file_path=str(file_path),
                        source=api_endpoint,
                        source_type=SourceType.API.value,
                        source_context=img_url,
                        file_size=len(content)
                    )
                    image_ids.append(img_data.generate_id())
                    
                except Exception as e:
                    logger.warning(f"Failed to download API image {img_url}: {e}")
            
            logger.info(f"Acquired {len(image_ids)} images from API: {api_endpoint}")
            
        except Exception as e:
            logger.error(f"Error acquiring from API {api_endpoint}: {e}")
        
        return image_ids
    
    def _parse_api_response(self, data: Any) -> List[str]:
        """Parse API response to extract image URLs"""
        urls = []
        
        if isinstance(data, list):
            for item in data:
                urls.extend(self._parse_api_response(item))
        elif isinstance(data, dict):
            # Common JSON paths for image URLs
            for key in ['image', 'image_url', 'url', 'src', 'thumbnail', 'images']:
                if key in data:
                    value = data[key]
                    if isinstance(value, str) and self._is_valid_url(value):
                        urls.append(value)
                    elif isinstance(value, list):
                        urls.extend(self._parse_api_response(value))
        elif isinstance(data, str) and self._is_valid_url(data):
            urls.append(data)
        
        return urls
    
    def _is_valid_url(self, url: str) -> bool:
        """Check if string is a valid URL"""
        try:
            result = urllib.parse.urlparse(url)
            return all([result.scheme, result.netloc])
        except:
            return False
    
    def _generate_thumbnail(self, image_path: str) -> Optional[str]:
        """Generate thumbnail for image"""
        try:
            with Image.open(image_path) as img:
                if img.mode in ('RGBA', 'P'):
                    img = img.convert('RGB')
                
                img.thumbnail((256, 256), Image.Resampling.LANCZOS)
                
                thumb_filename = f"thumb_{Path(image_path).name}"
                thumb_path = self.thumbnail_path / thumb_filename
                img.save(thumb_path, "JPEG", quality=85)
                
                return str(thumb_path)
                
        except Exception as e:
            logger.warning(f"Error generating thumbnail for {image_path}: {e}")
            return None
    
    def _classify_image_type(self, image_path: str) -> str:
        """Classify image type (placeholder for ML-based classification)"""
        filename = Path(image_path).stem.lower()
        
        # Simple keyword-based classification
        if any(kw in filename for kw in ['floor', 'plan', 'layout', 'blueprint']):
            return ImageType.FLOOR_PLAN.value
        elif any(kw in filename for kw in ['facility', 'equipment', 'machine', 'plant']):
            return ImageType.FACILITY.value
        elif any(kw in filename for kw in ['sign', 'label', 'symbol', 'mark', 'warning']):
            return ImageType.SIGN.value
        elif any(kw in filename for kw in ['standard', 'spec', 'diagram', 'schematic']):
            return ImageType.STANDARD.value
        elif any(kw in filename for kw in ['scene', 'photo', 'view', 'site']):
            return ImageType.SCENE.value
        
        return ImageType.UNKNOWN.value


# Utility functions
def get_image_info(image_id: str) -> Optional[Dict[str, Any]]:
    """Get image information by ID (placeholder - would query database)"""
    # This would query the database for image metadata
    return None


def list_acquired_images(source_type: Optional[str] = None) -> List[Dict[str, Any]]:
    """List all acquired images, optionally filtered by source type"""
    images = []
    
    try:
        for img_path in settings.images_raw.glob("*"):
            if img_path.suffix.lower() in ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp'):
                img_data = ImageData(
                    file_path=str(img_path),
                    source_type=source_type or "unknown"
                )
                images.append(img_data.to_dict())
    except Exception as e:
        logger.error(f"Error listing images: {e}")
    
    return images


# Export main classes
__all__ = [
    'ImageData',
    'ImageType',
    'SourceType',
    'DocumentImageExtractor',
    'WebImageCrawler',
    'FileUploadHandler',
    'ImageAcquisitionService',
    'get_image_info',
    'list_acquired_images'
]
